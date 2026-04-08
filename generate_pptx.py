from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to add slide
    def add_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        content = slide.placeholders[1]
        content.text = content_text

    # Helper for Title Slide
    def add_title_slide(main_text, sub_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = main_text
        subtitle.text = sub_text

    # --- SLIDE 1: TITLE ---
    add_title_slide("Sales Bot Platform", "Universal Conversational Commerce Engine\nVersión 2.0.0 | Multi-Producto")

    # --- SLIDE 2: OVERVIEW ---
    add_slide("Resumen Ejecutivo", 
    "• Agente de Ventas Autónomo diseñado para WhatsApp\n"
    "• 100% Agnóstico al Producto: Funciona para tecnología, ropa, comida, servicios\n"
    "• Configuración dinámica 'Zero-Code' vía variables de entorno\n"
    "• Soporte Multi-Tenant: Maneja múltiples marcas/tiendas\n"
    "• Estado: Production-Ready (98% code coverage)")

    # --- SLIDE 3: IA UNIVERSAL ---
    add_slide("1. Inteligencia Artificial Universal",
    "• Dual LLM Core: Soporte nativo para GPT-4 y Gemini 1.5 Pro\n"
    "• Contexto Dinámico: System prompts se adaptan automáticamente al tipo de tienda\n"
    "• Personalidad Configurable: Tono (formal/informal), idioma y emojis ajustables\n"
    "• Detección de Intención: Navegación, compra, soporte, reclamos")

    # --- SLIDE 4: GESTIÓN DE INVENTARIO ---
    add_slide("2. Gestión de Inventario Dinámica",
    "• Catálogo Flexible: Carga desde CSV con cualquier estructura de productos\n"
    "• Auto-Categorización: Detecta categorías (ej: Remeras, Herramientas) automáticamente\n"
    "• Búsqueda Fuzzy: Encuentra productos aunque el usuario tenga errores de tipeo\n"
    "• Sincronización: Integración bidireccional con Google Sheets en tiempo real")

    # --- SLIDE 5: ESTRATEGIAS DE VENTA ---
    add_slide("3. Estrategias de Venta Automáticas",
    "• Cross-Selling Dinámico: Reglas automáticas basadas en categorías complementarias\n"
    "• Upselling Inteligente: Detecta versiones superiores (precio/specs) y las ofrece\n"
    "• Bundles & Packs: Soporte para combos permanentes y promociones temporales\n"
    "• Recuperación: Seguimiento automático de carritos abandonados")

    # --- SLIDE 6: CONFIGURACIÓN ---
    add_slide("4. Configuración & Despliegue",
    "• Archivo .env Único: Define nombre, tipo de tienda y reglas en segundos\n"
    "• Feature Flags: Activa/desactiva módulos (Upsell, Fraud Detection) con flags\n"
    "• Docker Ready: Contenedores optimizados para producción\n"
    "• Multi-Cloud: Deploy fácil en Railway, Heroku, AWS o VPS")

    # --- SLIDE 7: INTEGRACIONES ---
    add_slide("5. Integraciones Enterprise",
    "• WhatsApp: Twilio y Meta Cloud API (Facebook)\n"
    "• Pagos: MercadoPago (Links automáticos y Webhooks)\n"
    "• Email: Confirmaciones de pedido HTML con branding\n"
    "• Monitoreo: Sentry para errores y Redis para caching")

    # Save
    output_file = "platform_features.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        create_presentation()
    except Exception as e:
        print(f"Error: {e}")
